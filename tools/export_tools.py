"""Rule-based export outline and lightweight artifact export tools."""

from __future__ import annotations

from collections.abc import Callable
from io import BytesIO
from typing import Any

from openpyxl import Workbook

from persistence.interfaces import ArtifactStore
from schemas._base import utc_now
from schemas.analysis_package import AnalysisPackage
from schemas.dashboard import (
    DashboardArtifactMetadata,
    DashboardFilter,
    DashboardLayout,
    DashboardSpec,
    DashboardWidget,
    DashboardWidgetType,
)
from schemas.query_result import QueryResult
from schemas.report import (
    ArtifactRef,
    ReportFormat,
    ReportOutline,
    ReportOutlineSection,
    ReportResult,
)

EXCEL_MIME_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
REPORT_MIME_TYPE = "text/markdown"
PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
DASHBOARD_MIME_TYPE = "application/json"
PPT_TABLE_ROW_LIMIT = 5
DASHBOARD_FILTER_VALUE_LIMIT = 20


class ExportToolError(RuntimeError):
    """Structured export tool failure with a stable error code."""

    def __init__(
        self,
        error_code: str,
        message: str,
        *,
        details: dict[str, Any] | None = None,
    ) -> None:
        super().__init__(message)
        self.error_code = error_code
        self.details = details or {}


def propose_report_outline(package: AnalysisPackage) -> ReportOutline:
    """Build a report outline preview from an existing analysis package."""

    return _build_outline(
        package=package,
        report_format=ReportFormat.REPORT,
        title=f"Analysis report: {package.question}",
        section_builders=(
            _summary_section,
            _insight_section,
            _chart_section,
            _data_section,
        ),
    )


def propose_ppt_outline(package: AnalysisPackage) -> ReportOutline:
    """Build a PPT outline preview from an existing analysis package."""

    return _build_outline(
        package=package,
        report_format=ReportFormat.PPT,
        title=f"Analysis deck: {package.question}",
        section_builders=(
            _ppt_opening_section,
            _insight_section,
            _chart_section,
            _next_steps_section,
        ),
    )


def propose_excel_export(package: AnalysisPackage) -> ReportOutline:
    """Build an Excel export preview from an existing analysis package."""

    return _build_outline(
        package=package,
        report_format=ReportFormat.EXCEL,
        title=f"Analysis workbook: {package.question}",
        section_builders=(
            _summary_sheet_section,
            _data_section,
            _chart_metadata_section,
        ),
    )


def propose_dashboard_outline(package: AnalysisPackage) -> ReportOutline:
    """Build a dashboard outline preview from an existing analysis package."""

    return _build_outline(
        package=package,
        report_format=ReportFormat.DASHBOARD,
        title=f"Analysis dashboard: {package.question}",
        section_builders=(
            _dashboard_overview_section,
            _dashboard_widget_section,
            _dashboard_filter_section,
        ),
    )


def export_report(
    outline: ReportOutline,
    artifact_store: ArtifactStore,
    package: AnalysisPackage | None = None,
) -> ReportResult:
    """Create a lightweight Markdown report artifact from an outline and package."""

    _validate_outline_format(
        outline=outline,
        tool_name="export_report",
        report_format=ReportFormat.REPORT,
    )
    metadata = _export_metadata(
        outline=outline,
        package=package,
        tool_name="export_report",
        report_format=ReportFormat.REPORT,
        report_type="markdown_report",
        mime_type=REPORT_MIME_TYPE,
        placeholder=False,
    )
    artifact_ref_value = artifact_store.save_artifact(
        content=_build_markdown_report(outline, package),
        metadata=metadata,
    )
    return _report_result(
        report_format=ReportFormat.REPORT,
        artifact_ref_value=artifact_ref_value,
        metadata=metadata,
        status="created",
    )


def generate_ppt(
    outline: ReportOutline,
    artifact_store: ArtifactStore,
    package: AnalysisPackage | None = None,
) -> ReportResult:
    """根据已确认的大纲和分析包生成真实 PPTX artifact。"""

    _validate_outline_format(
        outline=outline,
        tool_name="generate_ppt",
        report_format=ReportFormat.PPT,
    )
    metadata = _export_metadata(
        outline=outline,
        package=package,
        tool_name="generate_ppt",
        report_format=ReportFormat.PPT,
        report_type="pptx",
        mime_type=PPTX_MIME_TYPE,
        placeholder=False,
    )
    pptx_content = _build_pptx_deck_or_raise(outline, package)
    artifact_ref_value = artifact_store.save_artifact(
        content=pptx_content,
        metadata=metadata,
    )
    return _report_result(
        report_format=ReportFormat.PPT,
        artifact_ref_value=artifact_ref_value,
        metadata=metadata,
        status="created",
    )


def export_excel(
    outline: ReportOutline,
    artifact_store: ArtifactStore,
    package: AnalysisPackage | None = None,
) -> ReportResult:
    """Export AnalysisPackage query results into an XLSX artifact."""

    _validate_outline_format(
        outline=outline,
        tool_name="export_excel",
        report_format=ReportFormat.EXCEL,
    )
    metadata = _export_metadata(
        outline=outline,
        package=package,
        tool_name="export_excel",
        report_format=ReportFormat.EXCEL,
        report_type="excel",
        mime_type=EXCEL_MIME_TYPE,
        placeholder=False,
    )
    artifact_ref_value = artifact_store.save_artifact(
        content=_build_excel_workbook(outline, package),
        metadata=metadata,
    )
    return _report_result(
        report_format=ReportFormat.EXCEL,
        artifact_ref_value=artifact_ref_value,
        metadata=metadata,
        status="created",
    )


def generate_dashboard(
    outline: ReportOutline,
    artifact_store: ArtifactStore,
    package: AnalysisPackage | None = None,
) -> ReportResult:
    """Create a structured dashboard spec artifact without rendering a frontend."""

    _validate_outline_format(
        outline=outline,
        tool_name="generate_dashboard",
        report_format=ReportFormat.DASHBOARD,
    )
    dashboard_spec = _build_dashboard_spec(outline, package)
    dashboard_metadata = _dashboard_artifact_metadata(dashboard_spec)
    metadata = {
        **_export_metadata(
            outline=outline,
            package=package,
            tool_name="generate_dashboard",
            report_format=ReportFormat.DASHBOARD,
            report_type="dashboard_spec",
            mime_type=DASHBOARD_MIME_TYPE,
            placeholder=False,
        ),
        **dashboard_metadata.model_dump(mode="json"),
    }
    artifact_ref_value = artifact_store.save_artifact(
        content=dashboard_spec.model_dump(mode="json"),
        metadata=metadata,
    )
    return _report_result(
        report_format=ReportFormat.DASHBOARD,
        artifact_ref_value=artifact_ref_value,
        metadata=metadata,
        status="created",
    )


def _build_outline(
    *,
    package: AnalysisPackage,
    report_format: ReportFormat,
    title: str,
    section_builders: tuple[Callable[[AnalysisPackage], ReportOutlineSection], ...],
) -> ReportOutline:
    """Create an outline using deterministic sections from the package."""

    return ReportOutline(
        report_format=report_format,
        title=title,
        sections=[builder(package) for builder in section_builders],
        source_package_id=package.package_id,
        requires_confirmation=True,
    )


def _summary_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create a report summary section."""

    return ReportOutlineSection(
        title="Summary",
        points=[
            f"Question: {package.question}",
            _primary_insight_text(package),
        ],
    )


def _ppt_opening_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create a deck opening section."""

    return ReportOutlineSection(
        title="Opening",
        points=[
            "Title slide",
            f"Analysis question: {package.question}",
        ],
    )


def _summary_sheet_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create an Excel summary sheet section."""

    return ReportOutlineSection(
        title="Summary sheet",
        points=[
            f"Question: {package.question}",
            _primary_insight_text(package),
        ],
    )


def _dashboard_overview_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create a dashboard overview section."""

    return ReportOutlineSection(
        title="Overview",
        points=[
            f"Question: {package.question}",
            "Primary metric tiles from the analysis package",
        ],
    )


def _dashboard_widget_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create a dashboard widget preview section."""

    points = []
    if package.sql_result is not None:
        points.extend(
            [
                "Metric widget from query result summary",
                "Table widget references query result metadata without row payload",
            ]
        )
    if package.chart_spec is not None:
        artifact_ref = package.chart_spec.artifact_ref or "not generated"
        points.append(f"Chart widget references artifact: {artifact_ref}")
    if package.insights:
        points.append("Insight widget summarizes generated insights")
    return ReportOutlineSection(
        title="Widgets",
        points=points or ["Dashboard spec will use saved outline text."],
    )


def _dashboard_filter_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create a dashboard filter preview section."""

    if package.sql_result is None:
        points = ["No query-result filters are available."]
    else:
        points = [f"Filter candidate: {column.name}" for column in package.sql_result.columns]
    return ReportOutlineSection(title="Filters", points=points)


def _insight_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create an insight section from package insights."""

    insight_points = [insight.summary for insight in package.insights]
    return ReportOutlineSection(
        title="Insights",
        points=insight_points or ["No insight summary is available yet."],
    )


def _chart_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create a chart section from the chart spec reference."""

    if package.chart_spec is None:
        points = ["No chart spec is available yet."]
    else:
        points = [
            f"Chart type: {package.chart_spec.chart_type.value}",
            f"Title: {package.chart_spec.title or 'Untitled chart'}",
        ]
    return ReportOutlineSection(title="Visuals", points=points)


def _data_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create a data section from the bounded query result."""

    if package.sql_result is None:
        points = ["No query result is available yet."]
    else:
        column_names = ", ".join(column.name for column in package.sql_result.columns)
        points = [
            f"Rows: {package.sql_result.row_count}",
            f"Columns: {column_names or 'none'}",
        ]
    return ReportOutlineSection(title="Data", points=points)


def _chart_metadata_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create an Excel chart metadata sheet section."""

    if package.chart_spec is None:
        points = ["No chart metadata sheet is needed."]
    else:
        points = [
            f"Chart id: {package.chart_spec.chart_id}",
            f"Artifact ref: {package.chart_spec.artifact_ref or 'not generated'}",
        ]
    return ReportOutlineSection(title="Chart metadata", points=points)


def _next_steps_section(package: AnalysisPackage) -> ReportOutlineSection:
    """Create a deck next-steps section."""

    return ReportOutlineSection(
        title="Next steps",
        points=[
            "Confirm whether this outline should be exported.",
            f"Source package: {package.package_id}",
        ],
    )


def _primary_insight_text(package: AnalysisPackage) -> str:
    """Return the first insight summary or a deterministic fallback."""

    if package.insights:
        return package.insights[0].summary
    return "Analysis package contains no generated insight yet."


def _export_placeholder(
    *,
    outline: ReportOutline,
    artifact_store: ArtifactStore,
    tool_name: str,
    report_format: ReportFormat,
    package: AnalysisPackage | None = None,
) -> ReportResult:
    """Persist placeholder artifact metadata and return only references."""

    _validate_outline_format(outline=outline, tool_name=tool_name, report_format=report_format)

    metadata = _export_metadata(
        outline=outline,
        package=package,
        tool_name=tool_name,
        report_format=report_format,
        report_type=f"{report_format.value}_placeholder",
        mime_type="application/json",
        placeholder=True,
    )
    artifact_ref_value = artifact_store.save_artifact(
        content={
            "kind": "placeholder_export",
            "outline": outline.model_dump(mode="json"),
        },
        metadata=metadata,
    )
    return _report_result(
        report_format=report_format,
        artifact_ref_value=artifact_ref_value,
        metadata=metadata,
        status="placeholder_created",
    )


def _validate_outline_format(
    *,
    outline: ReportOutline,
    tool_name: str,
    report_format: ReportFormat,
) -> None:
    """Ensure an outline is compatible with the selected export tool."""

    if outline.report_format != report_format:
        raise ValueError(
            f"Outline format '{outline.report_format.value}' cannot be exported by '{tool_name}'."
        )


def _report_result(
    *,
    report_format: ReportFormat,
    artifact_ref_value: str,
    metadata: dict[str, Any],
    status: str,
) -> ReportResult:
    """Build a ReportResult from artifact metadata."""

    artifact = ArtifactRef(
        artifact_ref=artifact_ref_value,
        artifact_type=report_format,
        metadata=metadata,
    )
    return ReportResult(
        report_format=report_format,
        artifact_ref=artifact.artifact_ref,
        status=status,
        artifact=artifact,
    )


def _export_metadata(
    *,
    outline: ReportOutline,
    package: AnalysisPackage | None,
    tool_name: str,
    report_format: ReportFormat,
    report_type: str,
    mime_type: str,
    placeholder: bool,
) -> dict[str, Any]:
    """Build common artifact metadata for export outputs."""

    return {
        "tool_name": tool_name,
        "outline_id": outline.outline_id,
        "report_format": report_format.value,
        "report_type": report_type,
        "source_analysis_package_id": (
            package.package_id if package is not None else outline.source_package_id
        ),
        "created_at": utc_now().isoformat(),
        "mime_type": mime_type,
        "placeholder": placeholder,
    }


def _build_excel_workbook(outline: ReportOutline, package: AnalysisPackage | None) -> bytes:
    """Return XLSX bytes for a package's query result and outline metadata."""

    workbook = Workbook()
    summary_sheet = workbook.active
    summary_sheet.title = "Summary"
    summary_sheet.append(["Title", outline.title])
    summary_sheet.append(["Outline ID", outline.outline_id])
    summary_sheet.append(["Source package", _source_package_id(outline, package) or ""])
    if package is not None:
        summary_sheet.append(["Question", package.question])
        for index, insight in enumerate(package.insights, start=1):
            summary_sheet.append([f"Insight {index}", insight.summary])

    if package is not None and package.sql_result is not None:
        _append_query_result_sheet(workbook, package.sql_result)

    outline_sheet = workbook.create_sheet("Outline")
    outline_sheet.append(["Section", "Point"])
    for section in outline.sections:
        if section.points:
            for point in section.points:
                outline_sheet.append([section.title, point])
        else:
            outline_sheet.append([section.title, ""])

    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _append_query_result_sheet(workbook: Workbook, query_result: QueryResult) -> None:
    """Add bounded QueryResult rows to a workbook sheet."""

    worksheet = workbook.create_sheet("Query Result")
    column_names = [column.name for column in query_result.columns]
    if not column_names and query_result.rows:
        column_names = list(query_result.rows[0])
    worksheet.append(column_names)
    for row in query_result.rows:
        worksheet.append([row.get(column_name) for column_name in column_names])


def _build_markdown_report(outline: ReportOutline, package: AnalysisPackage | None) -> str:
    """Build a product-facing Markdown report without internal debug identifiers."""

    lines = [
        f"# {outline.title}",
        "",
        "## 数据源概览",
        "",
        *_datasource_overview_lines(package),
        "",
        "## 分析问题",
        "",
        (package.question if package is not None else "基于已保存大纲生成报告。"),
        "",
        "## 核心结论摘要",
        "",
        _executive_summary(package, outline),
        "",
        "## 主要发现",
        "",
        *_finding_lines(package, outline),
        "",
        "## 图表 / 表格说明",
        "",
        *_visual_and_table_lines(package),
        "",
        "## 数据限制",
        "",
        *_report_limitations(package),
        "",
        "## 后续建议",
        "",
        *_next_step_lines(package),
    ]
    return "\n".join(lines).strip() + "\n"


def _datasource_overview_lines(package: AnalysisPackage | None) -> list[str]:
    """Create a short datasource overview from bounded package metadata."""

    if package is None:
        return ["- 本报告基于已保存的报告大纲生成，未重新访问数据源。"]
    outputs: list[str] = []
    if package.sql_result is not None:
        outputs.append(f"查询结果 {package.sql_result.row_count} 行")
    if package.chart_spec is not None:
        outputs.append(f"{package.chart_spec.chart_type.value} 图表说明")
    if package.insights:
        outputs.append(f"{len(package.insights)} 条分析发现")
    summary = "、".join(outputs) if outputs else "已保存的分析上下文"
    return [f"- 本报告基于当前会话中已完成的分析结果生成，包含{summary}。"]


def _executive_summary(package: AnalysisPackage | None, outline: ReportOutline) -> str:
    """Return a concise executive summary in Chinese."""

    if package is not None and package.insights:
        return package.insights[0].summary
    for section in outline.sections:
        if section.points:
            return section.points[0]
    return "当前报告基于已保存的大纲生成，建议补充分析结果后继续完善。"


def _finding_lines(package: AnalysisPackage | None, outline: ReportOutline) -> list[str]:
    """Render findings without dumping internal package/debug identifiers."""

    if package is not None and package.insights:
        return [
            (
                f"{index}. **{_clean_report_text(insight.title)}**："
                f"{_clean_report_text(insight.summary)}"
            )
            for index, insight in enumerate(package.insights, start=1)
        ]
    lines: list[str] = []
    for section in outline.sections:
        for point in section.points:
            lines.append(f"- **{_clean_report_text(section.title)}**：{_clean_report_text(point)}")
    return lines or ["- 暂无可用发现，请先完成一次分析或开放探索。"]


def _visual_and_table_lines(package: AnalysisPackage | None) -> list[str]:
    """Describe chart/table outputs without exposing internal artifact ids."""

    lines: list[str] = []
    if package is not None and package.chart_spec is not None:
        lines.append(
            f"- 图表建议：使用 {package.chart_spec.chart_type.value} 图展示"
            f"“{package.chart_spec.title or '分析结果'}”。"
        )
    if package is not None and package.sql_result is not None:
        column_names = [column.name for column in package.sql_result.columns]
        clean_columns = ", ".join(_display_column_name(column) for column in column_names)
        table_result_line = (
            f"- 表格结果：查询返回 {package.sql_result.row_count} 行，"
            f"字段包括 {clean_columns or '暂无字段'}。"
        )
        lines.append(table_result_line)
        if not _query_result_uses_anonymous_columns(package.sql_result):
            lines.extend(["", _markdown_table(package.sql_result)])
    return lines or ["- 当前报告没有可展示的图表或表格结果。"]


def _report_limitations(package: AnalysisPackage | None) -> list[str]:
    """List data caveats in user-facing Chinese."""

    limitations = [
        "- 本报告复用当前会话中已经生成的分析结果，未重新分析或访问外部服务。",
    ]
    if package is not None and package.sql_result is not None:
        anonymous_columns = [
            column.name
            for column in package.sql_result.columns
            if _is_anonymous_report_column(column.name)
        ]
        for column_name in anonymous_columns:
            limitations.append(f"- {column_name} 字段缺少业务含义，仅作为未命名指标候选展示。")
    if package is None:
        limitations.append("- 当前只有报告大纲，没有可复用的查询结果或洞察。")
    return limitations


def _next_step_lines(package: AnalysisPackage | None) -> list[str]:
    """Suggest practical next actions."""

    suggestions = [
        "- 可以基于本报告继续生成 PPT、Excel 或 Dashboard。",
        "- 如果某个字段含义不清晰，建议先补充字段口径后再做正式决策。",
    ]
    if package is not None and package.insights:
        suggestions.append("- 可以选择某一条主要发现继续追问，做更细的维度拆解。")
    return suggestions


def _query_result_uses_anonymous_columns(query_result: QueryResult) -> bool:
    """Return whether a query result has generated anonymous measure names."""

    return any(_is_anonymous_report_column(column.name) for column in query_result.columns)


def _is_anonymous_report_column(column_name: str) -> bool:
    """Identify anonymous spreadsheet/generated metric names in report output."""

    lowered = column_name.lower()
    return (
        lowered.startswith("column_")
        or lowered.startswith("total_column_")
        or lowered.startswith("avg_column_")
        or lowered.startswith("unnamed")
    )


def _display_column_name(column_name: str) -> str:
    """Return a friendly display name for report table columns."""

    if _is_anonymous_report_column(column_name):
        return f"{column_name}（未命名指标）"
    return column_name


def _clean_report_text(value: str) -> str:
    """Remove known internal/debug markers from user-facing report text."""

    cleaned = str(value)
    for marker in ("Outline ID", "Source package", "Node name", "retry attempt"):
        cleaned = cleaned.replace(marker, "")
    return cleaned.strip()


def _markdown_table(query_result: QueryResult) -> str:
    """Render a small QueryResult as a Markdown table."""

    column_names = [column.name for column in query_result.columns]
    if not column_names and query_result.rows:
        column_names = list(query_result.rows[0])
    if not column_names:
        return "_No tabular result available._"
    header = "| " + " | ".join(column_names) + " |"
    separator = "| " + " | ".join("---" for _ in column_names) + " |"
    body = [
        "| " + " | ".join(str(row.get(column_name, "")) for column_name in column_names) + " |"
        for row in query_result.rows
    ]
    return "\n".join([header, separator, *body])


def _build_dashboard_spec(
    outline: ReportOutline,
    package: AnalysisPackage | None,
) -> DashboardSpec:
    """Build a dashboard spec that references artifacts instead of embedding bodies."""

    widgets = _dashboard_widgets(outline, package)
    return DashboardSpec(
        title=outline.title or "Analysis dashboard",
        source_package_id=_source_package_id(outline, package),
        question=package.question if package is not None else None,
        widgets=widgets,
        filters=_dashboard_filters(package),
    )


def _dashboard_widgets(
    outline: ReportOutline,
    package: AnalysisPackage | None,
) -> list[DashboardWidget]:
    """Create dashboard widgets from package outputs and saved outline text."""

    widgets: list[DashboardWidget] = []
    if package is not None and package.sql_result is not None:
        widgets.append(_metric_widget(package.sql_result, y=0))
        widgets.append(_table_widget(package.sql_result, y=3))
    if package is not None and package.chart_spec is not None:
        widgets.append(_chart_widget(package, y=0))
    if package is not None and package.insights:
        widgets.append(_insight_widget(package, y=6))
    if not widgets:
        widgets.append(_outline_text_widget(outline))
    return widgets


def _metric_widget(query_result: QueryResult, *, y: int) -> DashboardWidget:
    """Create a metric widget from the first numeric-looking result column."""

    metric_name = _metric_column_name(query_result)
    metric_value = None
    if metric_name is not None and query_result.rows:
        metric_value = query_result.rows[0].get(metric_name)
    return DashboardWidget(
        title="Primary metric",
        widget_type=DashboardWidgetType.METRIC,
        layout=DashboardLayout(x=0, y=y, w=3, h=2),
        query_result_id=query_result.result_id,
        metric_name=metric_name,
        value=metric_value,
        description="Metric value derived from the analysis package query result.",
        metadata={"row_count": query_result.row_count},
    )


def _chart_widget(package: AnalysisPackage, *, y: int) -> DashboardWidget:
    """Create a chart widget that references an existing chart artifact."""

    chart_spec = package.chart_spec
    if chart_spec is None:
        raise ValueError("ChartSpec is required for a chart dashboard widget.")
    return DashboardWidget(
        title=chart_spec.title or "Analysis chart",
        widget_type=DashboardWidgetType.CHART,
        layout=DashboardLayout(x=3, y=y, w=6, h=4),
        chart_type=chart_spec.chart_type,
        chart_artifact_ref=chart_spec.artifact_ref,
        description="Chart widget references the chart artifact; rendered content is not embedded.",
        metadata={"chart_id": chart_spec.chart_id},
    )


def _table_widget(query_result: QueryResult, *, y: int) -> DashboardWidget:
    """Create a table widget with schema metadata but no inline rows."""

    return DashboardWidget(
        title="Query result table",
        widget_type=DashboardWidgetType.TABLE,
        layout=DashboardLayout(x=0, y=y, w=6, h=3),
        query_result_id=query_result.result_id,
        description="Table widget points to the analysis query result and omits row payload.",
        metadata={
            "row_count": query_result.row_count,
            "columns": [column.name for column in query_result.columns],
        },
    )


def _insight_widget(package: AnalysisPackage, *, y: int) -> DashboardWidget:
    """Create an insight widget with concise summaries."""

    return DashboardWidget(
        title="Insights",
        widget_type=DashboardWidgetType.INSIGHT,
        layout=DashboardLayout(x=6, y=y, w=6, h=3),
        description="Key insight summaries from the analysis package.",
        metadata={
            "insights": [
                {
                    "title": insight.title,
                    "summary": insight.summary,
                }
                for insight in package.insights
            ]
        },
    )


def _outline_text_widget(outline: ReportOutline) -> DashboardWidget:
    """Create a fallback text widget when only a saved outline exists."""

    return DashboardWidget(
        title=outline.title or "Dashboard outline",
        widget_type=DashboardWidgetType.TEXT,
        layout=DashboardLayout(x=0, y=0, w=6, h=3),
        description="Saved dashboard outline without rerunning analysis.",
        metadata={
            "sections": [
                {
                    "title": section.title,
                    "points": section.points,
                }
                for section in outline.sections
            ]
        },
    )


def _dashboard_filters(package: AnalysisPackage | None) -> list[DashboardFilter]:
    """Create bounded filter metadata from query-result columns."""

    if package is None or package.sql_result is None:
        return []
    filters: list[DashboardFilter] = []
    for column in package.sql_result.columns:
        values = _filter_values(package.sql_result, column.name)
        if not values:
            continue
        filters.append(
            DashboardFilter(
                field=column.name,
                label=column.name.replace("_", " ").title(),
                values=values,
            )
        )
    return filters


def _dashboard_artifact_metadata(spec: DashboardSpec) -> DashboardArtifactMetadata:
    """Build dashboard artifact metadata without the full spec body."""

    return DashboardArtifactMetadata(
        dashboard_id=spec.dashboard_id,
        title=spec.title,
        source_analysis_package_id=spec.source_package_id,
        widget_count=len(spec.widgets),
        filter_count=len(spec.filters),
        chart_artifact_refs=_dashboard_chart_refs(spec),
        created_at=spec.created_at,
    )


def _dashboard_chart_refs(spec: DashboardSpec) -> list[str]:
    """Collect chart artifact references from dashboard widgets."""

    refs: list[str] = []
    for widget in spec.widgets:
        if widget.chart_artifact_ref is not None and widget.chart_artifact_ref not in refs:
            refs.append(widget.chart_artifact_ref)
    return refs


def _metric_column_name(query_result: QueryResult) -> str | None:
    """Select the first aggregate-looking column for the metric widget."""

    for column in query_result.columns:
        if column.name.startswith("total_") or column.name.startswith("sum_"):
            return column.name
    return query_result.columns[-1].name if query_result.columns else None


def _filter_values(query_result: QueryResult, column_name: str) -> list[Any]:
    """Collect bounded unique filter values without embedding full result rows."""

    values: list[Any] = []
    for row in query_result.rows:
        value = row.get(column_name)
        if value is not None and value not in values:
            values.append(value)
        if len(values) >= DASHBOARD_FILTER_VALUE_LIMIT:
            break
    return values


def _build_pptx_deck_or_raise(
    outline: ReportOutline,
    package: AnalysisPackage | None,
) -> bytes:
    """Build PPTX bytes and convert dependency/content issues into structured errors."""

    try:
        pptx_content = _build_pptx_deck(outline, package)
    except ModuleNotFoundError as exc:
        raise ExportToolError(
            "ppt_dependency_missing",
            "python-pptx is required to generate PPT exports.",
            details={"missing_module": exc.name or "pptx"},
        ) from exc
    except Exception as exc:
        raise ExportToolError(
            "ppt_generation_failed",
            "PPT export failed while building the pptx artifact.",
            details={"error_type": type(exc).__name__},
        ) from exc

    if not isinstance(pptx_content, bytes) or not pptx_content:
        raise ExportToolError(
            "ppt_invalid_artifact_content",
            "PPT export produced invalid artifact content.",
            details={"content_type": type(pptx_content).__name__},
        )
    return pptx_content


def _build_pptx_deck(
    outline: ReportOutline,
    package: AnalysisPackage | None,
) -> bytes:
    """为已确认的大纲构造结构化 PPTX 文件内容。"""

    from pptx import Presentation

    presentation = Presentation()
    _add_title_slide(presentation, outline, package)
    _add_bullet_slide(
        presentation,
        title="Summary",
        bullets=_summary_slide_bullets(outline, package),
    )
    _add_bullet_slide(
        presentation,
        title="SQL and Metric Result",
        bullets=_result_slide_bullets(package),
    )
    _add_bullet_slide(
        presentation,
        title="Insights",
        bullets=_insight_slide_bullets(package),
    )
    _add_bullet_slide(
        presentation,
        title="Chart Artifact Reference",
        bullets=_chart_slide_bullets(package),
    )
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _add_title_slide(
    presentation,
    outline: ReportOutline,
    package: AnalysisPackage | None,
) -> None:
    """添加标题页，并保证标题与副标题都有内容。"""

    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = outline.title or "Analysis deck"
    subtitle = slide.placeholders[1]
    subtitle.text = _title_slide_subtitle(outline, package)


def _add_bullet_slide(presentation, *, title: str, bullets: list[str]) -> None:
    """添加标题加要点页，并保证页面正文不为空。"""

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    safe_bullets = bullets or ["No content is available."]
    text_frame.text = safe_bullets[0]
    for bullet in safe_bullets[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0


def _title_slide_subtitle(outline: ReportOutline, package: AnalysisPackage | None) -> str:
    """返回标题页使用的非空副标题。"""

    if package is not None:
        return package.question
    if outline.source_package_id:
        return f"Source package: {outline.source_package_id}"
    return "Confirmed analysis export"


def _summary_slide_bullets(
    outline: ReportOutline,
    package: AnalysisPackage | None,
) -> list[str]:
    """生成 PPT 摘要页的简要要点。"""

    bullets = [
        f"Source package: {_source_package_id(outline, package) or 'unknown'}",
        f"Outline sections: {len(outline.sections)}",
    ]
    if package is not None:
        bullets.insert(0, f"Question: {package.question}")
        if package.insights:
            bullets.append(f"Primary insight: {package.insights[0].summary}")
    for section in outline.sections[:3]:
        if section.points:
            bullets.append(f"{section.title}: {section.points[0]}")
        else:
            bullets.append(section.title)
    return bullets


def _result_slide_bullets(package: AnalysisPackage | None) -> list[str]:
    """生成 SQL 与受限查询结果预览要点。"""

    if package is None or package.sql_result is None:
        return ["No SQL result is available for this export."]
    query_result = package.sql_result
    column_names = [column.name for column in query_result.columns]
    if not column_names and query_result.rows:
        column_names = list(query_result.rows[0])
    bullets = [
        f"SQL: {_truncate_text(query_result.sql, 180)}",
        f"Rows returned: {query_result.row_count}",
        f"Columns: {', '.join(column_names) or 'none'}",
    ]
    for row_index, row in enumerate(query_result.rows[:PPT_TABLE_ROW_LIMIT], start=1):
        bullets.append(f"Row {row_index}: {_format_row_for_ppt(row, column_names)}")
    if query_result.row_count > PPT_TABLE_ROW_LIMIT:
        bullets.append(f"Showing first {PPT_TABLE_ROW_LIMIT} rows only.")
    return bullets


def _insight_slide_bullets(package: AnalysisPackage | None) -> list[str]:
    """生成洞察页要点。"""

    if package is None or not package.insights:
        return ["No generated insights are available."]
    return [f"{insight.title}: {insight.summary}" for insight in package.insights]


def _chart_slide_bullets(package: AnalysisPackage | None) -> list[str]:
    """生成图表引用页要点，不嵌入真实图表内容。"""

    if package is None or package.chart_spec is None:
        return ["No chart artifact is available."]
    chart_spec = package.chart_spec
    return [
        f"Chart type: {chart_spec.chart_type.value}",
        f"Chart title: {chart_spec.title or 'Untitled chart'}",
        f"Chart artifact ref: {chart_spec.artifact_ref or 'not generated'}",
        f"Chart id: {chart_spec.chart_id}",
    ]


def _format_row_for_ppt(row: dict[str, Any], column_names: list[str]) -> str:
    """把一行查询结果压缩为适合 PPT 展示的预览文本。"""

    active_columns = column_names or list(row)
    return "; ".join(
        f"{column_name}={_truncate_text(str(row.get(column_name, '')), 40)}"
        for column_name in active_columns
    )


def _truncate_text(text: str, max_length: int) -> str:
    """限制写入 PPT 的文本长度。"""

    return text if len(text) <= max_length else f"{text[:max_length]}..."


def _source_package_id(outline: ReportOutline, package: AnalysisPackage | None) -> str | None:
    """Return the source package id without recomputing analysis."""

    return package.package_id if package is not None else outline.source_package_id


__all__ = [
    "export_excel",
    "export_report",
    "DASHBOARD_MIME_TYPE",
    "DASHBOARD_FILTER_VALUE_LIMIT",
    "EXCEL_MIME_TYPE",
    "ExportToolError",
    "PPTX_MIME_TYPE",
    "REPORT_MIME_TYPE",
    "PPT_TABLE_ROW_LIMIT",
    "generate_dashboard",
    "generate_ppt",
    "propose_dashboard_outline",
    "propose_excel_export",
    "propose_ppt_outline",
    "propose_report_outline",
]
