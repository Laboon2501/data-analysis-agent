"""Report graph 真实 PPTX 导出测试。"""

import json
from io import BytesIO

import pytest
from pptx import Presentation

from graphs.report_graph import build_report_graph
from nodes.runtime import NodeExecutionError
from persistence import InMemoryArtifactStore
from schemas import AgentCommand, AgentState, ChartSpec, ChartType, EventType, Insight, ReportFormat
from schemas.analysis_package import AnalysisPackage
from schemas.query_result import QueryColumn, QueryResult
from schemas.report import ReportOutline, ReportOutlineSection
from tools.export_tools import PPTX_MIME_TYPE


def test_report_graph_exports_ppt_fast_path_without_replanning() -> None:
    """ppt_confirm 应复用已保存大纲，并导出可读取的 PPTX 内容。"""

    artifact_store = InMemoryArtifactStore()
    package = _analysis_package()
    outline = _outline(package)

    state = _run_report_graph(
        AgentState(
            session_id="session-1",
            job_id="job-1",
            user_message="confirm ppt",
            command=AgentCommand.PPT_CONFIRM,
            analysis_package=package,
            report_outline=outline,
        ),
        artifact_store,
    )

    assert state.report_outline.outline_id == outline.outline_id
    assert state.report_outline.title == outline.title
    assert state.report_result is not None
    assert state.report_result.status == "created"

    record = artifact_store.get_artifact(state.report_result.artifact_ref)
    assert record is not None
    assert isinstance(record.content, bytes)
    assert record.metadata["mime_type"] == PPTX_MIME_TYPE
    assert record.metadata["report_type"] == "pptx"
    assert record.metadata["source_analysis_package_id"] == package.package_id
    assert record.metadata["placeholder"] is False

    presentation = Presentation(BytesIO(record.content))
    slide_texts = [_slide_text(slide) for slide in presentation.slides]

    assert len(presentation.slides) == 5
    assert all(_slide_title(slide) for slide in presentation.slides)
    assert all(_slide_text(slide) for slide in presentation.slides)
    assert "Saved ppt outline" in slide_texts[0]
    assert "Show monthly revenue trend" in slide_texts[0]
    assert "SQL and Metric Result" in slide_texts[2]
    assert "Revenue trend" in slide_texts[3]
    assert "artifact:chart-1" in slide_texts[4]


def test_report_graph_ppt_events_only_reference_artifact() -> None:
    """PPT 事件只应暴露 artifact 引用和元数据，不包含 deck 正文。"""

    artifact_store = InMemoryArtifactStore()
    package = _analysis_package()
    state = _run_report_graph(
        AgentState(
            session_id="session-1",
            job_id="job-1",
            user_message="confirm ppt",
            command=AgentCommand.PPT_CONFIRM,
            analysis_package=package,
            report_outline=_outline(package),
        ),
        artifact_store,
    )

    artifact_events = [
        event for event in state.events if event.event_type is EventType.ARTIFACT_REF
    ]

    assert len(artifact_events) == 1
    assert artifact_events[0].payload["artifact_ref"] == state.report_result.artifact_ref
    serialized_events = json.dumps(
        [event.model_dump(mode="json") for event in state.events],
        sort_keys=True,
    )
    assert "Row 1:" not in serialized_events
    assert "Revenue trend: Revenue increased across the period." not in serialized_events
    assert "artifact:chart-1" not in serialized_events


def test_report_graph_ppt_dependency_error_is_structured(monkeypatch) -> None:
    """PPT dependency failures should be visible as structured AgentState errors."""

    def missing_pptx(*_args, **_kwargs):
        raise ModuleNotFoundError("No module named 'pptx'", name="pptx")

    monkeypatch.setattr("tools.export_tools._build_pptx_deck", missing_pptx)
    package = _analysis_package()
    state = AgentState(
        session_id="session-1",
        job_id="job-1",
        user_message="confirm ppt",
        command=AgentCommand.PPT_CONFIRM,
        analysis_package=package,
        report_outline=_outline(package),
    )

    with pytest.raises(NodeExecutionError) as exc_info:
        build_report_graph(artifact_store=InMemoryArtifactStore()).invoke(state)

    failed_state = exc_info.value.state
    assert failed_state.errors[-1].code == "ppt_dependency_missing"
    assert failed_state.errors[-1].details["details"] == {"missing_module": "pptx"}
    assert failed_state.events[-1].event_type is EventType.ERROR
    assert failed_state.events[-1].payload["code"] == "ppt_dependency_missing"


def _run_report_graph(state: AgentState, artifact_store: InMemoryArtifactStore) -> AgentState:
    """运行 report graph 并校验返回状态。"""

    return AgentState.model_validate(
        build_report_graph(artifact_store=artifact_store).invoke(state)
    )


def _analysis_package() -> AnalysisPackage:
    """构造用于 PPT graph 导出测试的分析包。"""

    return AnalysisPackage(
        question="Show monthly revenue trend",
        sql_result=QueryResult(
            sql="SELECT month, SUM(revenue) AS total_revenue FROM orders GROUP BY month",
            columns=[
                QueryColumn(name="month", data_type="text"),
                QueryColumn(name="total_revenue", data_type="real"),
            ],
            rows=[
                {"month": "2026-01", "total_revenue": 100.0},
                {"month": "2026-02", "total_revenue": 210.0},
            ],
            row_count=2,
        ),
        chart_spec=ChartSpec(
            chart_type=ChartType.LINE,
            title="Monthly revenue",
            artifact_ref="artifact:chart-1",
        ),
        insights=[
            Insight(
                title="Revenue trend",
                summary="Revenue increased across the period.",
            )
        ],
    )


def _outline(package: AnalysisPackage) -> ReportOutline:
    """构造 confirm fast-path 使用的已保存 PPT 大纲。"""

    return ReportOutline(
        report_format=ReportFormat.PPT,
        title="Saved ppt outline",
        sections=[ReportOutlineSection(title="Summary", points=["Use this exact outline."])],
        source_package_id=package.package_id,
    )


def _slide_title(slide) -> str:
    """从 python-pptx slide 中提取标题。"""

    title = slide.shapes.title
    return "" if title is None else title.text.strip()


def _slide_text(slide) -> str:
    """从单页 slide 中提取全部文本。"""

    texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
    return "\n".join(texts)
