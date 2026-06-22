"""真实 PPTX artifact 导出测试。"""

from io import BytesIO

import pytest
from pptx import Presentation

from persistence import InMemoryArtifactStore
from schemas import ChartSpec, ChartType, Insight, ReportFormat
from schemas.analysis_package import AnalysisPackage
from schemas.query_result import QueryColumn, QueryResult
from schemas.report import ReportOutline, ReportOutlineSection
from tools.export_tools import (
    PPT_TABLE_ROW_LIMIT,
    PPTX_MIME_TYPE,
    ExportToolError,
    generate_ppt,
)


def test_generate_ppt_creates_readable_pptx_artifact() -> None:
    """PPT 导出应保存可读取、且包含必要页面的 deck。"""

    artifact_store = InMemoryArtifactStore()
    package = _analysis_package(row_count=PPT_TABLE_ROW_LIMIT + 2)
    outline = _outline(package)

    result = generate_ppt(outline, artifact_store, package)

    assert result.report_format is ReportFormat.PPT
    assert result.status == "created"
    record = artifact_store.get_artifact(result.artifact_ref)
    assert record is not None
    assert isinstance(record.content, bytes)
    assert record.metadata["mime_type"] == PPTX_MIME_TYPE
    assert record.metadata["report_type"] == "pptx"
    assert record.metadata["source_analysis_package_id"] == package.package_id
    assert record.metadata["placeholder"] is False

    presentation = Presentation(BytesIO(record.content))
    slide_texts = [_slide_text(slide) for slide in presentation.slides]

    assert len(presentation.slides) == 5
    assert all(_slide_title(slide) for slide in presentation.slides)
    assert all(_slide_text(slide) for slide in presentation.slides)
    assert "Revenue deck" in slide_texts[0]
    assert "Show monthly revenue trend" in slide_texts[0]
    assert "Summary" in slide_texts[1]
    assert "SQL and Metric Result" in slide_texts[2]
    assert "Insights" in slide_texts[3]
    assert "Chart Artifact Reference" in slide_texts[4]
    assert "artifact:chart-1" in slide_texts[4]


def test_generate_ppt_limits_table_rows() -> None:
    """PPT 结果页只应预览前 N 行数据。"""

    artifact_store = InMemoryArtifactStore()
    package = _analysis_package(row_count=PPT_TABLE_ROW_LIMIT + 2)

    result = generate_ppt(_outline(package), artifact_store, package)
    record = artifact_store.get_artifact(result.artifact_ref)
    presentation = Presentation(BytesIO(record.content))
    result_slide_text = _slide_text(presentation.slides[2])

    assert "Row 5:" in result_slide_text
    assert "Showing first 5 rows only." in result_slide_text
    assert "2026-06" not in result_slide_text
    assert "2026-07" not in result_slide_text


def test_generate_ppt_dependency_failure_is_structured(monkeypatch) -> None:
    """Missing python-pptx should use a stable export error code."""

    def missing_pptx(*_args, **_kwargs):
        raise ModuleNotFoundError("No module named 'pptx'", name="pptx")

    monkeypatch.setattr("tools.export_tools._build_pptx_deck", missing_pptx)

    with pytest.raises(ExportToolError) as exc_info:
        generate_ppt(_outline(_analysis_package(row_count=1)), InMemoryArtifactStore())

    assert exc_info.value.error_code == "ppt_dependency_missing"
    assert exc_info.value.details == {"missing_module": "pptx"}


def _analysis_package(row_count: int) -> AnalysisPackage:
    """构造行数超过 PPT 预览上限的分析包。"""

    rows = [
        {
            "month": f"2026-{index:02d}",
            "total_revenue": float(index * 100),
        }
        for index in range(1, row_count + 1)
    ]
    return AnalysisPackage(
        question="Show monthly revenue trend",
        sql_result=QueryResult(
            sql="SELECT month, SUM(revenue) AS total_revenue FROM orders GROUP BY month",
            columns=[
                QueryColumn(name="month", data_type="text"),
                QueryColumn(name="total_revenue", data_type="real"),
            ],
            rows=rows,
            row_count=len(rows),
        ),
        chart_spec=ChartSpec(
            chart_type=ChartType.LINE,
            title="Monthly revenue",
            artifact_ref="artifact:chart-1",
        ),
        insights=[
            Insight(
                title="Revenue trend",
                summary="Revenue increased across the period.",
            )
        ],
    )


def _outline(package: AnalysisPackage) -> ReportOutline:
    """构造绑定分析包的已保存 PPT 大纲。"""

    return ReportOutline(
        report_format=ReportFormat.PPT,
        title="Revenue deck",
        sections=[
            ReportOutlineSection(
                title="Summary",
                points=["Use this exact outline."],
            )
        ],
        source_package_id=package.package_id,
    )


def _slide_title(slide) -> str:
    """从 python-pptx slide 中提取标题。"""

    title = slide.shapes.title
    return "" if title is None else title.text.strip()


def _slide_text(slide) -> str:
    """从单页 slide 中提取全部文本。"""

    texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
    return "\n".join(texts)
